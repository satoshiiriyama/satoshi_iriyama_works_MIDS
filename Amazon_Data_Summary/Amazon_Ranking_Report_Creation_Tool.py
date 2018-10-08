import requests
import bs4
from datetime import datetime
from pytz import timezone
import pandas as pd
import numpy as np
import openpyxl
from openpyxl.styles.fonts import Font
from openpyxl.styles import PatternFill, Border, Side, Alignment, Protection, Font
from openpyxl import Workbook
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches


headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 6.1) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/41.0.2228.0 Safari/537.36',
    }

res = requests.get('https://www.amazon.com/gp/bestsellers/wireless/2407749011/ref=pd_zg_hrsr_wireless_1_2_last', headers=headers)
res.raise_for_status()
soup = bs4.BeautifulSoup(res.text, "html5lib")

res_2 = requests.get('https://www.amazon.com/Best-Sellers-Cell-Phones-Accessories-Unlocked/zgbs/wireless/2407749011/ref=zg_bs_pg_2?_encoding=UTF8&pg=2', headers=headers)
res_2.raise_for_status()
soup_2 = bs4.BeautifulSoup(res_2.text, "html5lib")


fmt_1 = "%Y-%m-%d"
fmt_2 = "%H:%M"
fmt_F = "%H%M"

now_pst = datetime.now(timezone('US/Pacific'))

now_date_fmt = now_pst.strftime(fmt_1)
now_time_fmt = now_pst.strftime(fmt_2)
now_time_fmt_F = now_pst.strftime(fmt_F)

product_codes = soup.select('.zg-item-immersion')
product_codes_2 = soup_2.select('.zg-item-immersion')

product_codes_merged = product_codes + product_codes_2 

product_list = []
for product_name in product_codes_merged:
    if str(product_name).find("p13n-sc-truncate") >= 0:
        product_name = product_name.select('.p13n-sc-truncate')[0].getText()
        product_name = product_name[13:-9]
        product_list.append(product_name)
    else:
        product_list.append("n/a")

vendor_list = []

for product in product_list:
    product_split = product.split(" ")
    vendor = product_split[0]
    if vendor == "Moto":
        vendor = "Motorola"
    vendor_list.append(vendor)

asin_list = []

for product_name in product_codes_merged:
    if str(product_name).find("a-size-small a-link-normal") >= 0:
        asin = product_name.select('.a-size-small.a-link-normal')   
        asin = str(asin)
        asin = asin[62:72]
        asin_list.append(asin)
    else:
        asin_list.append("n/a")

review_number_list = []

for product_name in product_codes_merged:
    if str(product_name).find("a-size-small a-link-normal") >= 0:
        review_number = product_name.select('.a-size-small.a-link-normal')[0].getText()
        review_number = review_number.replace(",", "")
        review_number = int(review_number)
        review_number_list.append(review_number)
    
    else:
        review_number_list.append(0)

price_list = []

for product_name in product_codes_merged:
    if str(product_name).find("p13n-sc-price") >= 0 and product_name.select('.p13n-sc-price')[0].getText().find("-") < 0:
        price = product_name.select('.p13n-sc-price')[0].getText()
        price = price.replace(",", "")
        price = price.replace("$", "")
        price = float(price)
        price_list.append(price)
    else:
        price_list.append(0)

star_rating_list = []

for product_name in product_codes_merged:
    if str(product_name).find("a-icon-alt") >= 0:
        star_rating = product_name.select('.a-icon-alt')[0].getText()
        star_rating = str(star_rating)
        star_rating = star_rating[0:3]
        star_rating = float(star_rating)
        star_rating_list.append(star_rating)
    else:
        star_rating_list.append("")

Rank_100_df = pd.DataFrame({"Rank": np.arange(1,101),
                            "Vendor": vendor_list,
                            "Product": product_list,
                            "ASIN": asin_list,
                            "Price": price_list,
                            "Star Rating": star_rating_list,
                            "Reviews": review_number_list,
                            "Date" : now_date_fmt,
                            "Time": now_time_fmt}) 

Rank_100_df.to_csv(f"Resources/Amazon_Top_100_Latest_Product_list.csv", index=False)

file_location = "Resources/180904_Amazon_Product_Spec_Master.csv"
amz_master_df = pd.read_csv(file_location)

final_df = Rank_100_df.merge(amz_master_df, on="ASIN", how="left")

final_df[""] = ""

final_df = final_df[["Rank", "Vendor", "Product_x", "ASIN", "Price", "Star Rating",
                     "Reviews", "Date", "Time", "", "Announced", "Size", "Resolution", "OS",
                     "Chipset", "Memory", "Rear Camera", "Front Camera", "Weight", "Battery", "URL"]]
final_df = final_df.rename(columns={"Product_x": "Product"})

final_df.to_excel(f"Output/{now_date_fmt}_{now_time_fmt_F}_Amazon_Rank100.xlsx", index=False)
final_df.to_csv(f"Output/{now_date_fmt}_{now_time_fmt_F}_Amazon_Rank100.csv", index=False)
final_df.to_csv(f"Resources/Latest_Amazon_Rank100_for_Carpet.csv", index=False)

wb =openpyxl.load_workbook(f"Output/{now_date_fmt}_{now_time_fmt_F}_Amazon_Rank100.xlsx")
sheet = wb['Sheet1']
sheet.column_dimensions['A'].width = 4
sheet.column_dimensions['B'].width = 9
sheet.column_dimensions['C'].width = 48
sheet.column_dimensions['D'].width = 7
sheet.column_dimensions['E'].width = 10
sheet.column_dimensions['F'].width = 8
sheet.column_dimensions['G'].width = 8
sheet.column_dimensions['H'].width = 10.5
sheet.column_dimensions['I'].width = 7
sheet.column_dimensions['J'].width = 1
sheet.column_dimensions['K'].width = 9
sheet.column_dimensions['L'].width = 7
sheet.column_dimensions['M'].width = 21
sheet.column_dimensions['N'].width = 11
sheet.column_dimensions['O'].width = 34
sheet.column_dimensions['P'].width = 12
sheet.column_dimensions['Q'].width = 8
sheet.column_dimensions['R'].width = 8
sheet.column_dimensions['S'].width = 11
sheet.column_dimensions['T'].width = 9
sheet.column_dimensions['U'].width = 10

sheet.freeze_panes = 'F2'

fill_1 = openpyxl.styles.PatternFill(patternType='solid', fgColor='c6e2ff', bgColor='c6e2ff')
fill_2 = openpyxl.styles.PatternFill(patternType='solid', fgColor='daedfe', bgColor='daedfe')
fill_3 = openpyxl.styles.PatternFill(patternType='solid', fgColor='00FFFFFF', bgColor='00FFFFFF')


for i in range(1,22):
    sheet.cell(row=1, column=i).fill = fill_1
    sheet.cell(row=1, column=i).alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
    
    for x in range(2, 104):
        if x % 2 == 1:
            sheet.cell(row=x, column=i).fill = fill_2

for r in range(1, 104):
    for c in range(1, 22):
        sheet.cell(row=r, column=c).border = Border(outline=True, 
                                                    left=Side(style='dotted', color='FF000000'),
                                                    right=Side(style='dotted', color='FF000000'), 
                                                    top=Side(style='dotted', color='FF000000'), 
                                                    bottom=Side(style='dotted', color='FF000000'))
  
sheet.sheet_view.showGridLines = False

for r in range(2, 104):
    if sheet.cell(row=r, column=5) != "":
        sheet.cell(row=r, column=5).number_format = '$  #,##0.00' 
    
    sheet.cell(row=r, column=6).number_format = '#.0'

    
for r in range(1, 104):
    sheet.cell(row=r, column=10).border = Border(outline=True, 
                                                    left=Side(style='dotted', color='FF000000'),
                                                    right=Side(style='dotted', color='FF000000'), 
                                                    top=Side(style=None), 
                                                    bottom=Side(style=None))
    sheet.cell(row=r, column=10).fill = fill_3

sheet.sheet_view.zoomScale = 85
    
wb.save(f"Output/{now_date_fmt}_{now_time_fmt_F}_Amazon_Rank100.xlsx")




fmt_1 = "%Y-%m-%d"
fmt_2 = "%H:%M"
fmt_F = "%H%M"

now_pst = datetime.now(timezone('US/Pacific'))

now_date_fmt = now_pst.strftime(fmt_1)
now_time_fmt = now_pst.strftime(fmt_2)
now_time_fmt_F = now_pst.strftime(fmt_F)

amazon_ranking = "Resources/Latest_Amazon_Rank100_for_Carpet.csv"
volume_path = "Resources/180914_Amazon_Vol_by_rank.csv"

rank_df = pd.read_csv(amazon_ranking, encoding="utf-8")

volume_df = pd.read_csv(volume_path, encoding="utf-8")

volume_series = pd.Series(volume_df["Volume"], name='Volume')
combined_df = pd.concat([rank_df, volume_series], axis=1)

combined_selected_df = combined_df[["Vendor", "Price", "Volume"]]

price_below_100 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_100_150 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_150_200 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_200_250 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_250_300 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_300_400 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_400_500 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_500_600 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_600_800 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_800_1000 = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}
price_1000_above = {"Samsung": 0, "Motorola": 0, "Apple": 0, "Huawei": 0, "LG": 0, "Nokia": 0, "Essential": 0, "BLU": 0, "Others": 0}

vendor_list = ["Samsung", "Motorola", "Apple", "Huawei", "LG", "Nokia", "Essential", "BLU", "Others"]

for index, row in combined_selected_df.iterrows():
    if row["Vendor"] == "Moto":
        combined_selected_df.at[index, "Vendor"] = "Motorola"
    elif row["Vendor"] == "Honor":
        combined_selected_df.at[index, "Vendor"] = "Huawei"
    elif row["Vendor"] not in vendor_list:
        combined_selected_df.at[index, "Vendor"] = "Others"


for index, row in combined_selected_df.iterrows():
    if row["Price"] <100:
        price_below_100[row["Vendor"]] += row["Volume"]
    elif row["Price"] <150:
        price_100_150[row["Vendor"]] += row["Volume"]
    elif row["Price"] <200:
        price_150_200[row["Vendor"]] += row["Volume"]
    elif row["Price"] <250:
        price_200_250[row["Vendor"]] += row["Volume"]
    elif row["Price"] <300:
        price_250_300[row["Vendor"]] += row["Volume"]
    elif row["Price"] <400:
        price_300_400[row["Vendor"]] += row["Volume"]
    elif row["Price"] <500:
        price_400_500[row["Vendor"]] += row["Volume"]
    elif row["Price"] <600:
        price_500_600[row["Vendor"]] += row["Volume"]
    elif row["Price"] <800:
        price_600_800[row["Vendor"]] += row["Volume"]
    elif row["Price"] <1000:
        price_800_1000[row["Vendor"]] += row["Volume"]
    else:
        price_1000_above[row["Vendor"]] += row["Volume"]        


Dicts_list = [price_below_100, price_100_150, price_150_200, price_200_250, price_250_300, price_300_400, price_400_500, 
         price_500_600, price_600_800, price_800_1000, price_1000_above]


final_df = pd.DataFrame(Dicts_list)
final_df = final_df.T
final_df = final_df.reindex(["Samsung", "Motorola", "Apple", "Huawei", "LG", "Nokia", "Essential", "BLU", "Others"])
final_df = final_df.rename(columns={0: "Below $100", 1: "$100-150", 2: "$150-200", 3: "$200-250", 4: "$250-300", 
                                    5: "$300-400", 6: "$400-500", 7: "$500-600", 8: "$600-800", 9: "$800-1000", 
                                    10: "Above $1000"})


final_df.loc["Total"] = final_df.sum()

final_df["Total"] = final_df[final_df.columns].sum(axis=1)

total_list = final_df.loc["Total"]

vendor_total_volume = final_df["Total"]

total_list_for_width = []

for i in range(11):
    percentage = total_list[i]/total_list[11] * 100
    percentage = round(percentage, 2)
    total_list_for_width.append(percentage)

xtick = [total_list_for_width[0]/2]
previous_tick = total_list_for_width[0]/2

for i in range(len(total_list_for_width)-1):
    tick_location = previous_tick + total_list_for_width[i]/2 + total_list_for_width[i+1]/2
    xtick.append(tick_location)
    previous_tick = tick_location

x_edge = xtick[10] + total_list_for_width[10]/2

final_df = final_df.drop(columns=["Total"])

samsung_list = final_df.loc["Samsung"]
motorola_list = final_df.loc["Motorola"]
apple_list = final_df.loc["Apple"]
huawei_list = final_df.loc["Huawei"]
lg_list = final_df.loc["LG"]
nokia_list = final_df.loc["Nokia"]
essential_list = final_df.loc["Essential"]
blu_list = final_df.loc["BLU"]
others_list = final_df.loc["Others"]

total = samsung_list + motorola_list + apple_list + huawei_list + lg_list + nokia_list + essential_list + blu_list + others_list

proportion_samsung = np.true_divide(samsung_list, total) * 100
proportion_motorola = np.true_divide(motorola_list, total) * 100
proportion_apple = np.true_divide(apple_list, total) * 100
proportion_huawei = np.true_divide(huawei_list, total) * 100
proportion_lg = np.true_divide(lg_list, total) * 100
proportion_nokia = np.true_divide(nokia_list, total) * 100
proportion_essential = np.true_divide(essential_list, total) * 100
proportion_blu = np.true_divide(blu_list, total) * 100
proportion_others = np.true_divide(others_list, total) * 100

price_range_list = ["Below $100", "$100-150", "$150-200", "$200-250", "$250-300", "$300-400", "$400-500", 
                    "$500-600", "$600-800", "$800-1000", "Above $1000"]

plt.figure(figsize=(20,8))

plt.bar(xtick, proportion_others, width=total_list_for_width, label='Others', color='silver', edgecolor="gray", bottom=proportion_samsung+proportion_motorola+proportion_apple+proportion_huawei+proportion_lg+proportion_nokia+proportion_essential+proportion_blu)
plt.bar(xtick, proportion_blu, width=total_list_for_width, label='Blu', color='deepskyblue', edgecolor="gray", bottom=proportion_samsung+proportion_motorola+proportion_apple+proportion_huawei+proportion_lg+proportion_nokia+proportion_essential)
plt.bar(xtick, proportion_essential, width=total_list_for_width, label='Essential', color='maroon', edgecolor="gray", bottom=proportion_samsung+proportion_motorola+proportion_apple+proportion_huawei+proportion_lg+proportion_nokia)
plt.bar(xtick, proportion_nokia, width=total_list_for_width, label='Nokia', color='orange', edgecolor="gray", bottom=proportion_samsung+proportion_motorola+proportion_apple+proportion_huawei+proportion_lg)
plt.bar(xtick, proportion_lg, width=total_list_for_width, label='LG', color='dimgray', edgecolor="gray", bottom=proportion_samsung+proportion_motorola+proportion_apple+proportion_huawei)
plt.bar(xtick, proportion_huawei, width=total_list_for_width, label='Huawei', color='yellow', edgecolor="gray", bottom=proportion_samsung+proportion_motorola+proportion_apple)
plt.bar(xtick, proportion_apple, width=total_list_for_width, label='Apple', color='magenta', edgecolor="gray", bottom=proportion_samsung+proportion_motorola)
plt.bar(xtick, proportion_motorola, width=total_list_for_width, label='Motorola', color='lightgreen', edgecolor="gray", bottom=proportion_samsung)
plt.bar(xtick, proportion_samsung, width=total_list_for_width, label='Samsung', color='royalblue', edgecolor="gray")

plt.xticks(xtick, price_range_list, rotation=90, fontsize=20)

plt.legend(loc=(1.02, 0.2), prop={'size': 20})

plt.xlim(0, x_edge)

plt.title(f"Amazon Top100 Carpet Chart {now_date_fmt} {now_time_fmt_F}")

plt.savefig(f"Output/{now_date_fmt}_{now_time_fmt_F}_Amazon_Rank100_Carpet.png", bbox_inches="tight")
plt.savefig("Resources/Amazon_Rank100_Carpet.png", bbox_inches="tight")

# plt.show()

sizes = total_list[0:11]
labels = price_range_list
colors = ["silver", "lightgreen", "limegreen", "lightblue", "dodgerblue", "yellow", "orange", "lightpink", "hotpink", "magenta", "red"]

plt.pie(sizes, labels=labels, autopct="%1.1f%%", colors=colors, shadow=False, startangle=90, counterclock=False)

plt.axis("equal")

plt.title(f"Amazon Top100 Price-band Proportion {now_date_fmt} {now_time_fmt_F}")

plt.savefig(f"Output/{now_date_fmt}_{now_time_fmt_F}_Amazon_Rank100_Price_Band.png", bbox_inches="tight")
plt.savefig("Resources/Amazon_Rank100_Price_Band.png", bbox_inches="tight")

# plt.show()

sizes = vendor_total_volume[0:9]
labels = vendor_list
colors = ["royalblue", "lightgreen", "magenta", "yellow", "dimgray", "orange", "maroon", "deepskyblue", "silver"]

plt.pie(sizes, labels=labels, autopct="%1.1f%%", colors=colors, shadow=False, startangle=90, counterclock=False)

plt.axis("equal")

plt.title(f"Amazon Top100 Vendor Volume Share {now_date_fmt} {now_time_fmt_F}")

plt.savefig(f"Output/{now_date_fmt}_{now_time_fmt_F}_Amazon_Rank100_Vendor_Share.png", bbox_inches="tight")
plt.savefig("Resources/Amazon_Rank100_Vendor_Share.png", bbox_inches="tight")

# plt.show()

prs = Presentation('Resources/AMZ_Master_template.pptx')
img_path = "Resources/Amazon_Rank100_Carpet.png"

slide = prs.slides[0]

x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(4.5)
slide.shapes.add_picture(img_path, x, y, cx, cy)

img_path = "Resources/Amazon_Rank100_Vendor_Share.png"
slide = prs.slides[1]

x, y, cx, cy = Inches(0.5), Inches(1.7), Inches(4), Inches(2.9)
slide.shapes.add_picture(img_path, x, y, cx, cy)

img_path = "Resources/Amazon_Rank100_Price_Band.png"
slide = prs.slides[1]

x, y, cx, cy = Inches(5.5), Inches(1.7), Inches(4), Inches(2.9)
slide.shapes.add_picture(img_path, x, y, cx, cy)

prs.save(f"Output/{now_date_fmt}_{now_time_fmt_F}_Amazon_Rank100_Data_Summary.pptx")